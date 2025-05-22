#!/usr/bin/env python3
"""
JarvisAI Document Processor
This service handles document processing for the JarvisAI RAG system:
1. Extracts text from various document formats (PDF, DOCX, TXT, etc.)
2. Splits documents into manageable chunks
3. Generates embeddings for document chunks
4. Stores document chunks and metadata in Neo4j knowledge graph
5. Creates relationships between related chunks
"""

import os
import sys
import json
import time
import logging
import hashlib
import shutil
from typing import List, Dict, Any, Optional, Tuple, Set
import re
import traceback
from pathlib import Path
from datetime import datetime
import tempfile
import subprocess
from concurrent.futures import ThreadPoolExecutor, as_completed
import signal
import queue
import threading

# Document processing libraries
import fitz  # PyMuPDF
import docx
from pptx import Presentation
import csv
from bs4 import BeautifulSoup
import chardet

# Neo4j connection
from neo4j import GraphDatabase, Driver

# Requests for API calls
import requests

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(sys.stdout),
        logging.FileHandler(os.path.join(os.path.dirname(__file__), "processor.log"))
    ]
)
logger = logging.getLogger('document_processor')

# Configuration
OLLAMA_BASE_URL = os.environ.get('OLLAMA_BASE_URL', 'http://localhost:11434')
EMBEDDING_MODEL = os.environ.get('EMBEDDING_MODEL', 'nomic-embed-text')
NEO4J_URI = os.environ.get('NEO4J_URI', 'bolt://localhost:7687')
NEO4J_USER = os.environ.get('NEO4J_USER', 'neo4j')
NEO4J_PASSWORD = os.environ.get('NEO4J_PASSWORD', 'password')
UPLOAD_DIR = os.environ.get('UPLOAD_DIR', '/app/uploads')
PROCESSED_DIR = os.environ.get('PROCESSED_DIR', '/app/processed_documents')
CHUNK_SIZE = int(os.environ.get('CHUNK_SIZE', '1024'))
CHUNK_OVERLAP = int(os.environ.get('CHUNK_OVERLAP', '200'))
POLL_INTERVAL = int(os.environ.get('POLL_INTERVAL', '5'))  # seconds
MAX_WORKERS = int(os.environ.get('MAX_WORKERS', '2'))

# Create directories if they don't exist
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(PROCESSED_DIR, exist_ok=True)

# File type extensions mapping
FILE_TYPES = {
    'pdf': ['.pdf'],
    'word': ['.docx', '.doc'],
    'powerpoint': ['.pptx', '.ppt'],
    'excel': ['.xlsx', '.xls', '.csv'],
    'text': ['.txt', '.md', '.rtf'],
    'html': ['.html', '.htm'],
    'code': ['.py', '.js', '.java', '.c', '.cpp', '.go', '.rs', '.ts', '.php', '.rb']
}

# Reverse mapping for extension to type
EXT_TO_TYPE = {}
for file_type, extensions in FILE_TYPES.items():
    for ext in extensions:
        EXT_TO_TYPE[ext] = file_type

# Global variables
processing_queue = queue.Queue()
neo4j_driver: Optional[Driver] = None
should_exit = threading.Event()


class DocumentChunk:
    """Represents a chunk of text from a document with metadata"""
    
    def __init__(
        self, 
        text: str, 
        doc_id: str,
        chunk_id: str,
        page_num: Optional[int] = None,
        section: Optional[str] = None,
        metadata: Optional[Dict[str, Any]] = None
    ):
        self.text = text
        self.doc_id = doc_id
        self.chunk_id = chunk_id
        self.page_num = page_num
        self.section = section
        self.metadata = metadata or {}
        self.embedding: Optional[List[float]] = None
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary for storage"""
        return {
            "text": self.text,
            "doc_id": self.doc_id,
            "chunk_id": self.chunk_id,
            "page_num": self.page_num,
            "section": self.section,
            "metadata": self.metadata,
            "embedding": self.embedding
        }
    
    @classmethod
    def from_dict(cls, data: Dict[str, Any]) -> 'DocumentChunk':
        """Create from dictionary"""
        chunk = cls(
            text=data["text"],
            doc_id=data["doc_id"],
            chunk_id=data["chunk_id"],
            page_num=data.get("page_num"),
            section=data.get("section"),
            metadata=data.get("metadata", {})
        )
        chunk.embedding = data.get("embedding")
        return chunk


def connect_to_neo4j() -> Optional[Driver]:
    """Connect to Neo4j database"""
    try:
        driver = GraphDatabase.driver(
            NEO4J_URI, 
            auth=(NEO4J_USER, NEO4J_PASSWORD)
        )
        # Test the connection
        with driver.session() as session:
            result = session.run("RETURN 'Connection successful' AS message")
            for record in result:
                logger.info(f"Neo4j connection: {record['message']}")
        
        # Initialize schema constraints if they don't exist
        with driver.session() as session:
            # Create constraints (only if they don't exist yet)
            session.run("""
                CREATE CONSTRAINT document_id IF NOT EXISTS
                FOR (d:Document) REQUIRE d.doc_id IS UNIQUE
            """)
            session.run("""
                CREATE CONSTRAINT chunk_id IF NOT EXISTS
                FOR (c:Chunk) REQUIRE c.chunk_id IS UNIQUE
            """)
            # Create vector index if it doesn't exist
            session.run("""
                CREATE VECTOR INDEX chunk_embedding IF NOT EXISTS
                FOR (c:Chunk) ON (c.embedding)
                OPTIONS {indexConfig: {metric: 'cosine'}}
            """)
            
        logger.info("Neo4j schema initialized")
        return driver
    except Exception as e:
        logger.error(f"Failed to connect to Neo4j: {e}")
        return None


def get_embedding(text: str) -> Optional[List[float]]:
    """Get embedding for text using Ollama API"""
    try:
        response = requests.post(
            f"{OLLAMA_BASE_URL}/api/embeddings",
            json={"model": EMBEDDING_MODEL, "prompt": text}
        )
        
        if response.status_code != 200:
            logger.error(f"Failed to get embedding: {response.status_code} - {response.text}")
            return None
        
        data = response.json()
        return data.get("embedding")
    except Exception as e:
        logger.error(f"Error getting embedding: {e}")
        return None


def extract_text_from_pdf(file_path: str) -> List[Tuple[int, str]]:
    """Extract text from PDF file with page numbers"""
    try:
        result = []
        with fitz.open(file_path) as doc:
            for page_num, page in enumerate(doc):
                text = page.get_text()
                if text.strip():  # Only add non-empty pages
                    result.append((page_num + 1, text))
        return result
    except Exception as e:
        logger.error(f"Error extracting text from PDF {file_path}: {e}")
        return []


def extract_text_from_docx(file_path: str) -> List[str]:
    """Extract text from DOCX file"""
    try:
        doc = docx.Document(file_path)
        return [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
    except Exception as e:
        logger.error(f"Error extracting text from DOCX {file_path}: {e}")
        return []


def extract_text_from_pptx(file_path: str) -> List[str]:
    """Extract text from PPTX file"""
    try:
        prs = Presentation(file_path)
        text_content = []
        
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            if slide_text:
                text_content.append("\n".join(slide_text))
        
        return text_content
    except Exception as e:
        logger.error(f"Error extracting text from PPTX {file_path}: {e}")
        return []


def extract_text_from_csv(file_path: str) -> str:
    """Extract text from CSV file"""
    try:
        # Detect encoding
        with open(file_path, 'rb') as f:
            result = chardet.detect(f.read())
            encoding = result['encoding']
        
        rows = []
        with open(file_path, 'r', encoding=encoding) as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(",".join(row))
        
        return "\n".join(rows)
    except Exception as e:
        logger.error(f"Error extracting text from CSV {file_path}: {e}")
        return ""


def extract_text_from_text_file(file_path: str) -> str:
    """Extract text from plain text file"""
    try:
        # Detect encoding
        with open(file_path, 'rb') as f:
            result = chardet.detect(f.read())
            encoding = result['encoding']
        
        with open(file_path, 'r', encoding=encoding) as f:
            return f.read()
    except Exception as e:
        logger.error(f"Error extracting text from text file {file_path}: {e}")
        return ""


def extract_text_from_html(file_path: str) -> str:
    """Extract text from HTML file"""
    try:
        # Detect encoding
        with open(file_path, 'rb') as f:
            result = chardet.detect(f.read())
            encoding = result['encoding']
        
        with open(file_path, 'r', encoding=encoding) as f:
            html_content = f.read()
        
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.extract()
        
        # Get text
        return soup.get_text(separator="\n")
    except Exception as e:
        logger.error(f"Error extracting text from HTML {file_path}: {e}")
        return ""


def convert_doc_to_docx(file_path: str) -> Optional[str]:
    """Convert DOC to DOCX using LibreOffice"""
    try:
        temp_dir = tempfile.mkdtemp()
        output_file = os.path.join(temp_dir, 'output.docx')
        
        # Run LibreOffice conversion
        cmd = [
            'libreoffice', 
            '--headless', 
            '--convert-to', 'docx', 
            '--outdir', temp_dir, 
            file_path
        ]
        
        process = subprocess.run(cmd, capture_output=True, text=True)
        if process.returncode != 0:
            logger.error(f"Error converting DOC to DOCX: {process.stderr}")
            return None
        
        # Find the output file
        for file in os.listdir(temp_dir):
            if file.endswith('.docx'):
                return os.path.join(temp_dir, file)
        
        return None
    except Exception as e:
        logger.error(f"Error converting DOC to DOCX: {e}")
        return None


def extract_text_from_document(file_path: str) -> List[Tuple[Optional[int], str]]:
    """Extract text from a document based on its file type"""
    file_ext = os.path.splitext(file_path)[1].lower()
    
    # PDF files
    if file_ext in FILE_TYPES['pdf']:
        return extract_text_from_pdf(file_path)
    
    # Word documents
    elif file_ext == '.docx':
        paragraphs = extract_text_from_docx(file_path)
        return [(None, para) for para in paragraphs]
    elif file_ext == '.doc':
        # Convert DOC to DOCX first
        docx_path = convert_doc_to_docx(file_path)
        if docx_path:
            paragraphs = extract_text_from_docx(docx_path)
            # Clean up temp file
            try:
                os.remove(docx_path)
            except:
                pass
            return [(None, para) for para in paragraphs]
        return []
    
    # PowerPoint
    elif file_ext in FILE_TYPES['powerpoint']:
        slides = extract_text_from_pptx(file_path)
        return [(i+1, slide) for i, slide in enumerate(slides)]  # Use slide number as page
    
    # CSV
    elif file_ext == '.csv':
        text = extract_text_from_csv(file_path)
        return [(None, text)]
    
    # Plain text files
    elif file_ext in FILE_TYPES['text']:
        text = extract_text_from_text_file(file_path)
        return [(None, text)]
    
    # HTML files
    elif file_ext in FILE_TYPES['html']:
        text = extract_text_from_html(file_path)
        return [(None, text)]
    
    # Code files
    elif file_ext in FILE_TYPES['code']:
        text = extract_text_from_text_file(file_path)
        return [(None, text)]
    
    # Unsupported format
    else:
        logger.warning(f"Unsupported file format: {file_ext}")
        return []


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> List[str]:
    """Split text into chunks with overlap"""
    if not text.strip():
        return []
    
    # Split by sentences to avoid cutting in the middle of a sentence
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks = []
    current_chunk = []
    current_size = 0
    
    for sentence in sentences:
        sentence_size = len(sentence)
        
        # If a single sentence is larger than chunk_size, split it by words
        if sentence_size > chunk_size:
            words = sentence.split()
            word_chunk = []
            word_chunk_size = 0
            
            for word in words:
                word_size = len(word) + 1  # Add 1 for space
                if word_chunk_size + word_size > chunk_size:
                    # Add current word chunk to chunks
                    if word_chunk:
                        chunks.append(" ".join(word_chunk))
                    word_chunk = [word]
                    word_chunk_size = word_size
                else:
                    word_chunk.append(word)
                    word_chunk_size += word_size
            
            # Add any remaining words as a chunk
            if word_chunk:
                chunks.append(" ".join(word_chunk))
            
        # Normal case - add sentence to current chunk
        elif current_size + sentence_size <= chunk_size:
            current_chunk.append(sentence)
            current_size += sentence_size
        else:
            # Current chunk is full, start a new one
            if current_chunk:
                chunks.append(" ".join(current_chunk))
            
            # Create overlap by copying some sentences from previous chunk
            overlap_start = max(0, len(current_chunk) - overlap // len(sentences[-1]) if sentences else 0)
            current_chunk = current_chunk[overlap_start:] + [sentence]
            current_size = sum(len(s) for s in current_chunk)
    
    # Add the last chunk if it's not empty
    if current_chunk:
        chunks.append(" ".join(current_chunk))
    
    return chunks


def process_document(file_path: str) -> Optional[str]:
    """Process a document and store it in Neo4j with embeddings"""
    if not os.path.exists(file_path):
        logger.error(f"File not found: {file_path}")
        return None
    
    try:
        # Generate document ID based on file content
        with open(file_path, 'rb') as f:
            file_hash = hashlib.sha256(f.read()).hexdigest()
        
        doc_id = f"doc_{file_hash[:16]}"
        file_name = os.path.basename(file_path)
        file_ext = os.path.splitext(file_name)[1].lower()
        file_type = EXT_TO_TYPE.get(file_ext, "unknown")
        
        # Extract document metadata
        file_stat = os.stat(file_path)
        creation_time = datetime.fromtimestamp(file_stat.st_ctime).isoformat()
        modified_time = datetime.fromtimestamp(file_stat.st_mtime).isoformat()
        file_size = file_stat.st_size
        
        # Document metadata
        doc_metadata = {
            "file_name": file_name,
            "file_type": file_type,
            "file_extension": file_ext,
            "file_size": file_size,
            "creation_time": creation_time,
            "modified_time": modified_time,
            "processed_time": datetime.now().isoformat()
        }
        
        # Extract text from document
        logger.info(f"Extracting text from {file_name}")
        text_parts = extract_text_from_document(file_path)
        
        if not text_parts:
            logger.warning(f"No text extracted from {file_name}")
            return None
        
        # Create document node in Neo4j
        if neo4j_driver:
            with neo4j_driver.session() as session:
                # Check if document already exists
                result = session.run(
                    "MATCH (d:Document {doc_id: $doc_id}) RETURN d",
                    doc_id=doc_id
                )
                
                if result.single():
                    logger.info(f"Document {doc_id} already exists in Neo4j")
                    # Delete existing document and its chunks
                    session.run(
                        """
                        MATCH (d:Document {doc_id: $doc_id})
                        OPTIONAL MATCH (d)-[:HAS_CHUNK]->(c:Chunk)
                        DETACH DELETE c, d
                        """,
                        doc_id=doc_id
                    )
                
                # Create document node
                session.run(
                    """
                    CREATE (d:Document {
                        doc_id: $doc_id,
                        file_name: $file_name,
                        file_type: $file_type,
                        file_extension: $file_ext,
                        file_size: $file_size,
                        creation_time: $creation_time,
                        modified_time: $modified_time,
                        processed_time: $processed_time
                    })
                    """,
                    doc_id=doc_id,
                    file_name=file_name,
                    file_type=file_type,
                    file_ext=file_ext,
                    file_size=file_size,
                    creation_time=creation_time,
                    modified_time=modified_time,
                    processed_time=doc_metadata["processed_time"]
                )
        
        # Process each text part
        all_chunks = []
        chunk_index = 0
        
        for page_num, text in text_parts:
            if not text.strip():
                continue
            
            # Split text into chunks
            text_chunks = chunk_text(text)
            
            for chunk_text in text_chunks:
                chunk_id = f"{doc_id}_chunk_{chunk_index}"
                chunk_index += 1
                
                # Create chunk object
                chunk = DocumentChunk(
                    text=chunk_text,
                    doc_id=doc_id,
                    chunk_id=chunk_id,
                    page_num=page_num,
                    metadata={
                        "file_name": file_name,
                        "file_type": file_type,
                        "chunk_index": chunk_index - 1
                    }
                )
                
                # Get embedding for the chunk
                chunk.embedding = get_embedding(chunk_text)
                
                if not chunk.embedding:
                    logger.warning(f"Failed to get embedding for chunk {chunk_id}")
                    continue
                
                all_chunks.append(chunk)
        
        # Store chunks in Neo4j
        if neo4j_driver and all_chunks:
            with neo4j_driver.session() as session:
                for chunk in all_chunks:
                    # Create chunk node and link to document
                    session.run(
                        """
                        MATCH (d:Document {doc_id: $doc_id})
                        CREATE (c:Chunk {
                            chunk_id: $chunk_id,
                            doc_id: $doc_id,
                            text: $text,
                            page_num: $page_num,
                            embedding: $embedding
                        })
                        CREATE (d)-[:HAS_CHUNK]->(c)
                        """,
                        doc_id=chunk.doc_id,
                        chunk_id=chunk.chunk_id,
                        text=chunk.text,
                        page_num=chunk.page_num,
                        embedding=chunk.embedding
                    )
            
            # Create relationships between chunks based on similarity
            with neo4j_driver.session() as session:
                # Create SIMILAR_TO relationships between chunks in the same document
                session.run(
                    """
                    MATCH (c1:Chunk {doc_id: $doc_id})
                    MATCH (c2:Chunk {doc_id: $doc_id})
                    WHERE c1.chunk_id < c2.chunk_id
                    WITH c1, c2, gds.similarity.cosine(c1.embedding, c2.embedding) AS similarity
                    WHERE similarity > 0.7
                    CREATE (c1)-[:SIMILAR_TO {score: similarity}]->(c2)
                    """,
                    doc_id=doc_id
                )
                
                # Create NEXT relationships between consecutive chunks
                session.run(
                    """
                    MATCH (d:Document {doc_id: $doc_id})-[:HAS_CHUNK]->(c:Chunk)
                    WITH d, c ORDER BY c.chunk_id
                    WITH d, collect(c) AS chunks
                    UNWIND range(0, size(chunks)-2) AS i
                    WITH chunks[i] AS c1, chunks[i+1] AS c2
                    CREATE (c1)-[:NEXT]->(c2)
                    """,
                    doc_id=doc_id
                )
        
        # Move the processed file to the processed directory
        processed_path = os.path.join(PROCESSED_DIR, file_name)
        shutil.move(file_path, processed_path)
        
        logger.info(f"Document {file_name} processed successfully with {len(all_chunks)} chunks")
        return doc_id
    
    except Exception as e:
        logger.error(f"Error processing document {file_path}: {e}")
        logger.error(traceback.format_exc())
        return None


def document_processor_worker():
    """Worker function to process documents from the queue"""
    while not should_exit.is_set():
        try:
            # Get a file from the queue with timeout
            file_path = processing_queue.get(timeout=1)
            
            logger.info(f"Processing file: {file_path}")
            doc_id = process_document(file_path)
            
            if doc_id:
                logger.info(f"Successfully processed document: {file_path} -> {doc_id}")
            else:
                logger.error(f"Failed to process document: {file_path}")
            
            # Mark the task as done
            processing_queue.task_done()
            
        except queue.Empty:
            # No files in the queue, just continue
            continue
        except Exception as e:
            logger.error(f"Error in document processing worker: {e}")
            logger.error(traceback.format_exc())
            # If we had a file path, mark it as done
            if 'file_path' in locals():
                processing_queue.task_done()


def check_for_files():
    """Check for new files in the upload directory"""
    try:
        for filename in os.listdir(UPLOAD_DIR):
            file_path = os.path.join(UPLOAD_DIR, filename)
            
            # Skip directories
            if os.path.isdir(file_path):
                continue
            
            # Check if file extension is supported
            file_ext = os.path.splitext(filename)[1].lower()
            if file_ext not in EXT_TO_TYPE:
                logger.warning(f"Unsupported file type: {file_ext} for file {filename}")
                # Move to processed with unsupported_ prefix
                unsupported_path = os.path.join(PROCESSED_DIR, f"unsupported_{filename}")
                shutil.move(file_path, unsupported_path)
                continue
            
            # Add file to processing queue
            logger.info(f"Adding file to processing queue: {filename}")
            processing_queue.put(file_path)
    
    except Exception as e:
        logger.error(f"Error checking for files: {e}")


def signal_handler(sig, frame):
    """Handle termination signals"""
    logger.info("Received termination signal, shutting down...")
    should_exit.set()
    
    # Close Neo4j connection
    if neo4j_driver:
        neo4j_driver.close()
    
    sys.exit(0)


def create_requirements_file():
    """Create requirements.txt file if it doesn't exist"""
    requirements_path = os.path.join(os.path.dirname(__file__), "requirements.txt")
    
    if not os.path.exists(requirements_path):
        with open(requirements_path, "w") as f:
            f.write("\n".join([
                "pymupdf==1.23.6",
                "python-docx==0.8.11",
                "python-pptx==0.6.21",
                "beautifulsoup4==4.12.2",
                "chardet==5.2.0",
                "neo4j==5.15.0",
                "requests==2.31.0"
            ]))
        logger.info(f"Created requirements.txt at {requirements_path}")


def main():
    """Main function to start the document processor service"""
    global neo4j_driver
    
    # Create requirements file if needed
    create_requirements_file()
    
    # Set up signal handlers
    signal.signal(signal.SIGINT, signal_handler)
    signal.signal(signal.SIGTERM, signal_handler)
    
    # Connect to Neo4j
    logger.info("Connecting to Neo4j...")
    retries = 0
    max_retries = 12  # Try for 2 minutes (12 * 10 seconds)
    
    while retries < max_retries and not should_exit.is_set():
        neo4j_driver = connect_to_neo4j()
        if neo4j_driver:
            break
        
        retries += 1
        logger.info(f"Retrying Neo4j connection in 10 seconds... ({retries}/{max_retries})")
        time.sleep(10)
    
    if not neo4j_driver and not should_exit.is_set():
        logger.error("Failed to connect to Neo4j after multiple attempts. Exiting.")
        return
    
    # Start worker threads
    logger.info(f"Starting {MAX_WORKERS} worker threads...")
    workers = []
    for _ in range(MAX_WORKERS):
        worker = threading.Thread(target=document_processor_worker, daemon=True)
        worker.start()
        workers.append(worker)
    
    # Main loop to check for new files
    logger.info(f"Monitoring {UPLOAD_DIR} for new files...")
    
    while not should_exit.is_set():
        check_for_files()
        time.sleep(POLL_INTERVAL)
    
    # Wait for all workers to finish
    logger.info("Waiting for worker threads to finish...")
    for worker in workers:
        worker.join(timeout=5)
    
    logger.info("Document processor service stopped")


if __name__ == "__main__":
    main()